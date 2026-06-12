import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# Theme Colors
COLOR_NAVY = RGBColor(0x00, 0x1F, 0x3F)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GOLD = RGBColor(0xFF, 0x8C, 0x00)
COLOR_LIGHT_GOLD = RGBColor(0xFF, 0xF8, 0xF0)
COLOR_LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
COLOR_BORDER_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
COLOR_TEXT_DARK = RGBColor(0x1A, 0x36, 0x5D)
COLOR_TEXT_MUTED = RGBColor(0x4A, 0x55, 0x68)

# Functional Component Colors
COLOR_BLUE_LAYER = RGBColor(0x31, 0x82, 0xCE)
COLOR_PURPLE_LAYER = RGBColor(0x80, 0x5A, 0xD5)
COLOR_ORANGE_LAYER = RGBColor(0xDD, 0x6B, 0x20)
COLOR_RED_LAYER = RGBColor(0xC5, 0x30, 0x30)
COLOR_GREEN_LAYER = RGBColor(0x2F, 0x85, 0x5A)

FONT_NAME = "Segoe UI"

def add_header_bar(slide, custom_tagline=None):
    """Adds a standardized top navigation/header bar across corporate slides."""
    # Header background
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_NAVY
    header_shape.line.color.rgb = COLOR_GOLD
    header_shape.line.width = Pt(2)
    
    # Left Platform Title
    tx_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(6), Inches(0.5))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AUTONOMOUS AI DBA OPERATIONS PLATFORM"
    p.font.name = FONT_NAME
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Right Tagline / Tracker
    tx_box_r = slide.shapes.add_textbox(Inches(7.333), Inches(0.15), Inches(5.6), Inches(0.5))
    tf_r = tx_box_r.text_frame
    p_r = tf_r.paragraphs[0]
    p_r.alignment = PP_ALIGN.RIGHT
    p_r.text = custom_tagline if custom_tagline else "Powered by OpenAI | Azure-Ready"
    p_r.font.name = FONT_NAME
    p_r.font.size = Pt(14)
    p_r.font.color.rgb = COLOR_GOLD

def add_slide_title(slide, text):
    """Creates a consistent styled presentation slide header."""
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY

def add_footer(slide, current_slide, total_slides=12):
    """Appends structural document tracking strings to slide footers."""
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Presenter: Shami Gaffar  |  Date: June 2026  |  Audience: Management & Technical Reviewers"
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_TEXT_MUTED
    
    p_r = tf.add_paragraph()
    p_r.alignment = PP_ALIGN.RIGHT
    p_r.text = f"Slide {current_slide} of {total_slides}"
    p_r.font.name = FONT_NAME
    p_r.font.size = Pt(10)
    p_r.font.color.rgb = COLOR_TEXT_MUTED

def add_image_placeholder(slide, left, top, width, height, text):
    """Generates clean structural layout boxes mimicking system/UI mockups."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    shape.line.color.rgb = COLOR_BORDER_GRAY
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f"[IMAGE PLACEHOLDER]\n{text}"
    p.font.name = FONT_NAME
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_MUTED

# ==============================================================================
# SLIDE 1: TITLE SLIDE
# ==============================================================================
slide1 = prs.slides.add_slide(blank_layout)
bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = COLOR_NAVY

# Background Image Box Visual Layer
add_image_placeholder(slide1, Inches(1), Inches(1), Inches(11.333), Inches(5.5), "Futuristic Database Server Room / AI Neural Network Background Visualization")

# Re-overlay clear text cards over placeholder
title_box = slide1.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(1.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "Autonomous AI DBA Operations Platform"
p.font.name = FONT_NAME
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.text = "Agentic AI for Microsoft SQL Server Administration"
p2.font.name = FONT_NAME
p2.font.size = Pt(28)
p2.font.color.rgb = COLOR_WHITE

meta_box = slide1.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
tf_m = meta_box.text_frame
p_m = tf_m.paragraphs[0]
p_m.alignment = PP_ALIGN.CENTER
p_m.text = "Presenter: Shami Gaffar | Date: June 2026\nProject Assignment Demo and Assessment Review"
p_m.font.name = FONT_NAME
p_m.font.size = Pt(18)
p_m.font.color.rgb = COLOR_WHITE

badge_box = slide1.shapes.add_textbox(Inches(1), Inches(6.0), Inches(11.333), Inches(0.8))
tf_b = badge_box.text_frame
p_b = tf_b.paragraphs[0]
p_b.alignment = PP_ALIGN.CENTER
p_b.text = "[ Powered by OpenAI ]   •   [ MCP Orchestration ]   •   [ Azure-Ready ]"
p_b.font.name = FONT_NAME
p_b.font.size = Pt(16)
p_b.font.bold = True
p_b.font.color.rgb = COLOR_GOLD

# Minimal Bottom Row Icon Blocks
add_image_placeholder(slide1, Inches(5.0), Inches(6.7), Inches(0.8), Inches(0.4), "SQL Server")
add_image_placeholder(slide1, Inches(6.0), Inches(6.7), Inches(0.8), Inches(0.4), "Python")
add_image_placeholder(slide1, Inches(7.0), Inches(6.7), Inches(0.8), Inches(0.4), "Azure")

# ==============================================================================
# SLIDE 2: PROBLEM STATEMENT
# ==============================================================================
slide2 = prs.slides.add_slide(blank_layout)
add_header_bar(slide2, "The Challenge Layer")
add_slide_title(slide2, "The Problem — Manual DBA Operations Are Broken")

# Bullets Left
bullet_box = slide2.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.5), Inches(3.2))
tf_bu = bullet_box.text_frame
tf_bu.word_wrap = True
bullets = [
    "Manual DBA operations are slow, reactive, and inherently error-prone",
    "DBAs spend hours daily on repetitive health monitoring tasks",
    "SQL Server failures often go undetected until production & user impact",
    "No unified AI-assisted triage and governance layer exists today",
    "Business Impact: Costly downtime, data loss risks, compliance failures"
]
for i, txt in enumerate(bullets):
    p = tf_bu.paragraphs[0] if i == 0 else tf_bu.add_paragraph()
    p.text = "• " + txt
    p.font.name = FONT_NAME
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT_DARK
    p.space_after = Pt(12)

# Callout Box Bottom
callout = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.4), Inches(6.5), Inches(1.3))
callout.fill.solid()
callout.fill.fore_color.rgb = COLOR_LIGHT_GOLD
callout.line.color.rgb = COLOR_GOLD
tf_c = callout.text_frame
tf_c.word_wrap = True
p_c = tf_c.paragraphs[0]
p_c.alignment = PP_ALIGN.CENTER
p_c.text = '"Can we build an AI Agent that monitors SQL Server, detects issues, and recommends governed remediation — automatically?"'
p_c.font.name = FONT_NAME
p_c.font.size = Pt(16)
p_c.font.bold = True
p_c.font.color.rgb = COLOR_NAVY

# Visual Right
add_image_placeholder(slide2, Inches(7.6), Inches(2.0), Inches(5.1), Inches(4.7), "The Problem Today: Stressed DBA managing multi-screen incident alarms and broken infrastructure logs")
add_footer(slide2, 2)

# ==============================================================================
# SLIDE 3: SOLUTION OVERVIEW
# ==============================================================================
slide3 = prs.slides.add_slide(blank_layout)
add_header_bar(slide3, "Core Innovations")
add_slide_title(slide3, "Solution — Autonomous AI DBA Operations Platform")

features = [
    {"title": "1. Monitor", "desc": "9 real-time SQL Server health checks running continuously across system vectors."},
    {"title": "2. Analyze", "desc": "OpenAI GPT-4o-mini powered Root Cause Analysis (RCA) and dynamic tuning engine."},
    {"title": "3. Govern", "desc": "RBAC-validated, human approval-gated automated remediation architectural layer."},
    {"title": "4. Automate", "desc": "Azure Monitor + Azure Automation adapter emitting secure Email & Teams triage alerts."}
]

grid_positions = [
    (Inches(0.6), Inches(2.0)),  # Top Left
    (Inches(4.8), Inches(2.0)),  # Top Right
    (Inches(0.6), Inches(4.5)),  # Bottom Left
    (Inches(4.8), Inches(4.5))   # Bottom Right
]

for idx, feat in enumerate(features):
    pos = grid_positions[idx]
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos[0], pos[1], Inches(3.9), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    box.line.color.rgb = COLOR_BORDER_GRAY
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = feat["title"]
    p.font.name = FONT_NAME
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    
    p2 = tf.add_paragraph()
    p2.text = feat["desc"]
    p2.font.name = FONT_NAME
    p2.font.size = Pt(15)
    p2.font.color.rgb = COLOR_TEXT_DARK
    p2.space_before = Pt(8)

# Center-Right Orchestration Flow Chart Area
add_image_placeholder(slide3, Inches(9.0), Inches(2.0), Inches(3.8), Inches(4.7), "Central Orchestration Model: Connecting telemetry endpoints securely to processing clusters")
add_footer(slide3, 3)

# ==============================================================================
# SLIDE 4: ARCHITECTURE DIAGRAM
# ==============================================================================
slide4 = prs.slides.add_slide(blank_layout)
add_header_bar(slide4, "Pipeline Topography")
add_slide_title(slide4, "Platform Architecture — 12-Layer Agentic Flow")

layers = [
    ("1. SQL Server Backend", COLOR_BLUE_LAYER),
    ("2. Telemetry Collectors", COLOR_BLUE_LAYER),
    ("3. Azure Log Analytics", COLOR_BLUE_LAYER),
    ("4. MCP Orchestration Layer", COLOR_PURPLE_LAYER),
    ("5. Agent Engine Core", COLOR_PURPLE_LAYER),
    ("6. OpenAI LLM Gate", COLOR_PURPLE_LAYER),
    ("7. RCA Diagnostics Suite", COLOR_ORANGE_LAYER),
    ("8. Advisory Engine", COLOR_ORANGE_LAYER),
    ("9. Authorization Gateway", COLOR_RED_LAYER),
    ("10. Automation Runner", COLOR_RED_LAYER),
    ("11. Notification Relays", COLOR_GREEN_LAYER),
    ("12. Logging & Audit Logs", COLOR_GREEN_LAYER)
]

# Render as an intelligent 4x3 structural pipeline matrix block
box_w, box_h = Inches(2.8), Inches(1.1)
x_start, y_start = Inches(0.6), Inches(2.0)
gap_x, gap_y = Inches(0.3), Inches(0.4)

for index, (layer_name, color) in enumerate(layers):
    row = index // 4
    col = index % 4
    
    cx = x_start + col * (box_w + gap_x)
    cy = y_start + row * (box_h + gap_y)
    
    shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(3)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = layer_name
    p.font.name = FONT_NAME
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK

# Quick Bottom Ecosystem Indicators
add_image_placeholder(slide4, Inches(0.6), Inches(6.4), Inches(2.5), Inches(0.4), "[Icon] SQL Server")
add_image_placeholder(slide4, Inches(3.8), Inches(6.4), Inches(2.5), Inches(0.4), "[Icon] Azure Hub")
add_image_placeholder(slide4, Inches(7.0), Inches(6.4), Inches(2.5), Inches(0.4), "[Icon] OpenAI Engine")

add_footer(slide4, 4)

# ==============================================================================
# SLIDE 5: MCP TOOL REGISTRY
# ==============================================================================
slide5 = prs.slides.add_slide(blank_layout)
add_header_bar(slide5, "Interface Capability Registry")
add_slide_title(slide5, "MCP Tool Orchestration Layer — 13 Registered DBA Tools")

# Column 1 Left Card Container
c1_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.3))
c1_box.fill.solid()
c1_box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
c1_box.line.color.rgb = COLOR_BORDER_GRAY
tf_c1 = c1_box.text_frame
p = tf_c1.paragraphs[0]
p.text = "Telemetry & Monitoring Capabilities"
p.font.name = FONT_NAME
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_NAVY
p.space_after = Pt(10)

m_tools = [
    "CHECK_CPU → CPU Engine Load Monitor",
    "CHECK_BLOCKING → Lock/Blocking Intersector",
    "CHECK_LONG_RUNNING_QUERIES → Thread Runtime",
    "CHECK_FAILED_JOBS → Automation Job Monitor",
    "CHECK_BACKUP_STATUS → Disaster Recovery Evaluator",
    "CHECK_DATABASE_SPACE → Allocation Volume Map",
    "CHECK_INDEX_FRAGMENTATION → Fragmentation Profiler",
    "CHECK_STATISTICS_HEALTH → Optimizer Stat Scan"
]
for t in m_tools:
    p = tf_c1.add_paragraph()
    p.text = f"⚙️ {t}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_DARK

# Column 2 Right Card Container
c2_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.3))
c2_box.fill.solid()
c2_box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
c2_box.line.color.rgb = COLOR_BORDER_GRAY
tf_c2 = c2_box.text_frame
p = tf_c2.paragraphs[0]
p.text = "Governance & Administration Capabilities"
p.font.name = FONT_NAME
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_NAVY
p.space_after = Pt(10)

g_tools = [
    "VALIDATE_RBAC_PERMISSION → Credential Guard",
    "REQUEST_FAILED_JOB_RESTART_APPROVAL → Approval Hub",
    "GENERATE_DAILY_HEALTH_REPORT → Analytic Engine",
    "GENERATE_PERFORMANCE_TUNING_REPORT → T-SQL Tuning",
    "SEND_TO_AZURE_MONITOR → External Cloud Syncer",
    "CREATE_AZURE_AUTOMATION_RUNBOOK_REQUEST → Runbooks"
]
for t in g_tools:
    p = tf_c2.add_paragraph()
    p.text = f"🛡️ {t}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_DARK

# Footer Badge
badge = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4))
badge.fill.solid()
badge.fill.fore_color.rgb = COLOR_NAVY
p_b = badge.text_frame.paragraphs[0]
p_b.alignment = PP_ALIGN.CENTER
p_b.text = "All 13 tools registered, orchestrated, and executing safely via standard MCP Server API Architecture"
p_b.font.size = Pt(12)
p_b.font.bold = True
p_b.font.color.rgb = COLOR_GOLD

add_footer(slide5, 5)

# ==============================================================================
# SLIDE 6: AGENTIC WORKFLOW
# ==============================================================================
slide6 = prs.slides.add_slide(blank_layout)
add_header_bar(slide6, "Operational Sequence Map")
add_slide_title(slide6, "Agentic DBA Workflow — 7-Step End to End Execution")

steps = [
    ("Step 0", "Environment Validation", "Verify connection bounds"),
    ("Step 1", "SQL Monitoring", "CPU, Block, Query runtime"),
    ("Step 2", "Extended Checks", "Jobs, Backup, Space profiles"),
    ("Step 3", "RBAC Validation", "Verify processing context"),
    ("Step 4", "AI Deep Analysis", "GPT-4o-mini RCA engine"),
    ("Step 5", "Governance Gate", "Construct remediation path"),
    ("Step 6", "Notification Hub", "Emit logging & team feeds")
]

for idx, (st_id, st_name, st_desc) in enumerate(steps):
    sx = Inches(0.5) + idx * Inches(1.75)
    
    card = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, Inches(2.2), Inches(1.6), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    card.line.color.rgb = COLOR_NAVY
    
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = st_id
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_GOLD
    
    p2 = tf.add_paragraph()
    p2.text = st_name
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_DARK
    
    p3 = tf.add_paragraph()
    p3.text = st_desc
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_TEXT_MUTED

# Governance Evaluation Unit
gov_box = slide6.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(4.5), Inches(4.8), Inches(4.333), Inches(1.6))
gov_box.fill.solid()
gov_box.fill.fore_color.rgb = COLOR_LIGHT_GOLD
gov_box.line.color.rgb = COLOR_GOLD
p_gov = gov_box.text_frame.paragraphs[0]
p_gov.alignment = PP_ALIGN.CENTER
p_gov.text = "Lead DBA Decision Matrix\nApprove or Reject Action?"
p_gov.font.size = Pt(13)
p_gov.font.bold = True
p_gov.font.color.rgb = COLOR_NAVY

# Execution Tracks
t1 = slide6.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(3.8), Inches(0.8))
t1.text_frame.word_wrap = True
p_t1 = t1.text_frame.paragraphs[0]
p_t1.text = "✔ IF APPROVED:\nAzure Automation Engine Triggered"
p_t1.font.color.rgb = COLOR_GREEN_LAYER
p_t1.font.bold = True

t2 = slide6.shapes.add_textbox(Inches(9.0), Inches(5.3), Inches(3.8), Inches(0.8))
t2.text_frame.word_wrap = True
p_t2 = t2.text_frame.paragraphs[0]
p_t2.text = "❌ IF REJECTED:\nRemediation Halted & Securely Audited"
p_t2.font.color.rgb = COLOR_RED_LAYER
p_t2.font.bold = True

add_footer(slide6, 6)

# ==============================================================================
# SLIDE 7: LIVE DEMO SCREENSHOT SLIDE
# ==============================================================================
slide7 = prs.slides.add_slide(blank_layout)
add_header_bar(slide7, "Interface Control Room")
add_slide_title(slide7, "Live Demo — Real SQL Server Data in Action")

demos = [
    ("Dashboard Analytics — KPI Cards", "FastAPI + Bootstrap dashboard showing global system operational stability.", Inches(0.6), Inches(2.0)),
    ("Monitoring Layer — Live Engine Blocking", "Session 51 blocked by Session 64, LCK_M_X active context, 827 seconds runtime.", Inches(6.9), Inches(2.0)),
    ("Approvals Control Terminal — Open Actions", "AI-provisioned structural pipeline authorization frame for failed task step restart.", Inches(0.6), Inches(4.5)),
    ("Audit Matrix Logging — Immutable Path", "Granular historical process tracking showing telemetry parameters and timestamps.", Inches(6.9), Inches(4.5))
]

for title, desc, px, py in demos:
    add_image_placeholder(slide7, px, py, Inches(5.8), Inches(1.6), title)
    
    lbl = slide7.shapes.add_textbox(px, py + Inches(1.65), Inches(5.8), Inches(0.6))
    tf = lbl.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_DARK

# Live Banner Accent
live_badge = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.7), Inches(2.1), Inches(0.8), Inches(0.3))
live_badge.fill.solid()
live_badge.fill.fore_color.rgb = COLOR_RED_LAYER
p_lv = live_badge.text_frame.paragraphs[0]
p_lv.alignment = PP_ALIGN.CENTER
p_lv.text = "LIVE DATA"
p_lv.font.size = Pt(10)
p_lv.font.bold = True
p_lv.font.color.rgb = COLOR_WHITE

note_box = slide7.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4))
p_n = note_box.text_frame.paragraphs[0]
p_n.alignment = PP_ALIGN.CENTER
p_n.text = "🔒 Operational state pulled live from AdventureWorks2019 SQL Server production replica instances — zero mock data layers utilized."
p_n.font.size = Pt(12)
p_n.font.italic = True
p_n.font.color.rgb = COLOR_TEXT_MUTED

add_footer(slide7, 7)

# ==============================================================================
# SLIDE 8: AI ANALYSIS SAMPLE
# ==============================================================================
slide8 = prs.slides.add_slide(blank_layout)
add_header_bar(slide8, "Cognitive Diagnostic Window")
add_slide_title(slide8, "AI-Powered Root Cause Analysis — Real Output")

# Left Box: Telemetry Dump
left_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.6))
left_box.fill.solid()
left_box.fill.fore_color.rgb = COLOR_NAVY
tf_l = left_box.text_frame
tf_l.word_wrap = True
p_l = tf_l.paragraphs[0]
p_l.text = "SQL MONITORING DATA STREAM (INPUT TO AI):\n" \
           "=========================================\n\n" \
           "• Active Blocking Chain: Session 51 <-- Session 64\n" \
           "• Wait Resource Type: LCK_M_X (Exclusive Lock)\n" \
           "• Cumulative Block Duration: 827 Seconds\n" \
           "• Head Block Query: UPDATE Person.Person SET ...\n" \
           "• Job Engine Flag: 2 Failures inside [DemoFailedJobs]\n" \
           "• Disaster Recovery State: BACKUP AGE CRITICAL - AdventureWorks2019"
p_l.font.name = "Courier New"
p_l.font.size = Pt(13)
p_l.font.color.rgb = RGBColor(0, 255, 0)

# Right Box: AI Response Evaluation
right_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.6))
right_box.fill.solid()
right_box.fill.fore_color.rgb = COLOR_LIGHT_GOLD
right_box.line.color.rgb = COLOR_GOLD
tf_r = right_box.text_frame
tf_r.word_wrap = True

p_r = tf_r.paragraphs[0]
p_r.text = "AI AGENT EVALUATION REPORT (OUTPUT):\n"
p_r.font.size = Pt(16)
p_r.font.bold = True
p_r.font.color.rgb = COLOR_GOLD

report_text = [
    "Incident Vector: High impact active exclusive transaction blocking chain.",
    "Root Cause Identification: Long running transactional updates lack index resolution coverage on Person.Person, bottlenecking down-stream system operations.",
    "System Risk Evaluation: SEVERE RISK LAYER",
    "Prescribed Action Items Plan:",
    "  1. Terminate processing head blocker session 64 (Approval Required)",
    "  2. Pass target update block variables to query tuning advisors",
    "  3. Force generation of targeted differential transaction backup logs",
    "  4. Rebuild underlying nonclustered indexing infrastructure structures"
]
for line in report_text:
    p = tf_r.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_DARK
    p.space_after = Pt(4)

add_footer(slide8, 8)

# ==============================================================================
# SLIDE 9: APPROVAL WORKFLOW DEMO
# ==============================================================================
slide9 = prs.slides.add_slide(blank_layout)
add_header_bar(slide9, "Human-In-The-Loop Governance")
add_slide_title(slide9, "Approval-Gated Remediation — Human in the Loop")

proc_steps = [
    ("Step 1", "AI Detects Threat", "System identifies critical lock timeouts or database engine fault patterns."),
    ("Step 2", "Ticket Generation", "Secure tracking Token generated, mapped to state database, and posted to UI panels."),
    ("Step 3", "DBA Evaluation", "Operations Engineer examines analytics context and commits an Approve/Deny response."),
    ("Step 4", "Secure Execution", "If permitted, cloud infrastructure execute scripts. Else, logs record explicit drop paths.")
]

for idx, (st, name, dsc) in enumerate(proc_steps):
    px = Inches(0.6) + idx * Inches(3.1)
    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, Inches(2.2), Inches(2.8), Inches(3.2))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    card.line.color.rgb = COLOR_BORDER_GRAY
    slide9.shapes._spTree.remove(card._element) # fix context binding
    slide9.shapes._spTree.append(card._element)
    
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.text = st
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_GOLD
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = name
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_NAVY
    p2.space_before = Pt(6)
    
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.text = dsc
    p3.font.size = Pt(12)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_before = Pt(12)

# Global Principle Callout
banner = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.8), Inches(12.133), Inches(0.9))
banner.fill.solid()
banner.fill.fore_color.rgb = COLOR_NAVY
p_bn = banner.text_frame.paragraphs[0]
p_bn.alignment = PP_ALIGN.CENTER
p_bn.text = "CORE COMPLIANCE RULE: Zero modification execution processes are permitted to run within database engine containers without manual operational verification bounds."
p_bn.font.size = Pt(13)
p_bn.font.bold = True
p_bn.font.color.rgb = COLOR_WHITE

add_footer(slide9, 9)

# ==============================================================================
# SLIDE 10: REAL vs SIMULATED
# ==============================================================================
slide10 = prs.slides.add_slide(blank_layout)
add_header_bar(slide10, "Architectural Verification Mapping")
add_slide_title(slide10, "Transparency — What is Real vs Simulated")

# Left Column: Real
r_header = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0), Inches(5.8), Inches(0.6))
r_header.fill.solid()
r_header.fill.fore_color.rgb = COLOR_GREEN_LAYER
r_header.line.fill.background()
p = r_header.text_frame.paragraphs[0]
p.text = "✔ REAL — Working Live Production Code"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE

r_body = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.6), Inches(5.8), Inches(4.0))
r_body.fill.solid()
r_body.fill.fore_color.rgb = COLOR_LIGHT_GRAY
r_body.line.color.rgb = COLOR_BORDER_GRAY
tf_rb = r_body.text_frame
tf_rb.word_wrap = True
real_items = [
    "SQL Server connections and active analytical indexing queries",
    "Real-time CPU tracking, database locks, and performance scans",
    "Automated backup metadata retrieval and configuration tracking",
    "AI Deep Root Cause Analysis processing models (GPT-4o-mini)",
    "Secure Human-in-the-loop authorization state workflows",
    "FastAPI responsive custom engine administration dashboard panels",
    " Granular security RBAC user processing boundaries validation"
]
for item in real_items:
    p = tf_rb.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_DARK
    p.space_after = Pt(4)

# Right Column: Simulated
s_header = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.8), Inches(0.6))
s_header.fill.solid()
s_header.fill.fore_color.rgb = COLOR_ORANGE_LAYER
s_header.line.fill.background()
p = s_header.text_frame.paragraphs[0]
p.text = "⚠ SIMULATED — Cloud Adapter Integration Ready"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE

s_body = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.6), Inches(5.8), Inches(4.0))
s_body.fill.solid()
s_body.fill.fore_color.rgb = COLOR_LIGHT_GRAY
s_body.line.color.rgb = COLOR_BORDER_GRAY
tf_sb = s_body.text_frame
tf_sb.word_wrap = True
sim_items = [
    "Azure Monitor Sync: Cloud adapters coded; environment pending config",
    "Azure Automation: Runbook integration complete; API endpoints bypassed",
    "SMTP Alerting System: Message formatting complete; mail relays skipped",
    "Microsoft Teams Webhooks: Webhook connectors validated using console logs",
    "SQL Agent Sub-System: Local execution environments fallback to MSDB emulation layers"
]
for item in sim_items:
    p = tf_sb.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_DARK
    p.space_after = Pt(6)

add_footer(slide10, 10)

# ==============================================================================
# SLIDE 11: DELIVERABLES AND SCORE
# ==============================================================================
slide11 = prs.slides.add_slide(blank_layout)
add_header_bar(slide11, "Project Milestones Evaluation")
add_slide_title(slide11, "Deliverables Status and Assessment Score")

# Left Side Table Representation
left_card = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0), Inches(6.5), Inches(4.7))
left_card.fill.solid()
left_card.fill.fore_color.rgb = COLOR_LIGHT_GRAY
left_card.line.color.rgb = COLOR_BORDER_GRAY
tf_table = left_card.text_frame
tf_table.word_wrap = True

p_th = tf_table.paragraphs[0]
p_th.text = "Deliverable Matrix Mapping               |  Status"
p_th.font.size = Pt(16)
p_th.font.bold = True
p_th.font.color.rgb = COLOR_NAVY
p_th.space_after = Pt(10)

delivs = [
    ("Problem Proposal Specifications Document ", "✔ FULL PASS"),
    ("12-Layer System Architecture Typology Map  ", "✔ FULL PASS"),
    ("Core Functional Agentic AI Framework Core   ", "✔ FULL PASS"),
    ("MCP Tools Registry Orchestration Server    ", "✔ FULL PASS"),
    ("Performance Auditing & Analysis Suite      ", "✔ FULL PASS"),
    ("Security Access Matrix (RBAC + Audit Logs) ", "✔ FULL PASS"),
    ("Disaster Recovery Telemetry Engine         ", "⚡ PARTIAL"),
    ("Automated Processing Infrastructure Hooks  ", "✔ FULL PASS"),
    ("FastAPI Operational Web Dashboard Panels  ", "✔ FULL PASS")
]
for name, status in delivs:
    p = tf_table.add_paragraph()
    p.text = f"• {name} -> {status}"
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.font.color.rgb = COLOR_TEXT_DARK

# Right Side Score Card
right_card = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(2.0), Inches(5.1), Inches(4.7))
right_card.fill.solid()
right_card.fill.fore_color.rgb = COLOR_NAVY
right_card.line.fill.background()
tf_rc = right_card.text_frame
tf_rc.word_wrap = True

p_sc = tf_rc.paragraphs[0]
p_sc.alignment = PP_ALIGN.CENTER
p_sc.text = "AGGREGATE ASSESSMENT SCORE"
p_sc.font.size = Pt(16)
p_sc.font.color.rgb = COLOR_WHITE

p_num = tf_rc.add_paragraph()
p_num.alignment = PP_ALIGN.CENTER
p_num.text = "84 / 100"
p_num.font.size = Pt(64)
p_num.font.bold = True
p_num.font.color.rgb = COLOR_GOLD

p_badge = tf_rc.add_paragraph()
p_badge.alignment = PP_ALIGN.CENTER
p_badge.text = "[ STATUS: STRONG PASS — DEMO READY ]"
p_badge.font.size = Pt(16)
p_badge.font.bold = True
p_badge.font.color.rgb = RGBColor(0, 255, 0)

add_image_placeholder(slide11, Inches(9.4), Inches(5.2), Inches(1.5), Inches(1.2), "Trophy Icon")

add_footer(slide11, 11)

# ==============================================================================
# SLIDE 12: NEXT PHASE AND CLOSING
# ==============================================================================
slide12 = prs.slides.add_slide(blank_layout)
add_header_bar(slide12, "Strategic Enterprise Horizon")
add_slide_title(slide12, "Next Phase Roadmap and Closing Summary")

# Left Section: Roadmap Cards
roadmap_box = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.7))
roadmap_box.fill.solid()
roadmap_box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
roadmap_box.line.color.rgb = COLOR_BORDER_GRAY
tf_rm = roadmap_box.text_frame
tf_rm.word_wrap = True

p = tf_rm.paragraphs[0]
p.text = "Strategic Multi-Phase Roadmap"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_NAVY
p.space_after = Pt(10)

phases = [
    ("Phase 2 — Production Scaling (Weeks 4-6)", "Bind live Azure native components, wire real credentials vaults, connect enterprise notification webhooks, scale to SQL standard engines."),
    ("Phase 3 — Deep Context Learning (Weeks 7-9)", "Implement RAG vector data pipelines over infrastructure change histories to achieve predictive database load profiling matrices."),
    ("Phase 4 — Controlled Autonomous Evolution", "Enable dynamic closed-loop script adjustments bound inside strictly constrained configuration spaces mapped by executive signing modules.")
]
for p_title, p_desc in phases:
    p = tf_rm.add_paragraph()
    p.text = f"🚀 {p_title}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    
    p = tf_rm.add_paragraph()
    p.text = p_desc
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_DARK
    p.space_after = Pt(6)

# Right Section: Summary Theme Statement Card
summary_box = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.8), Inches(3.4))
summary_box.fill.solid()
summary_box.fill.fore_color.rgb = COLOR_LIGHT_GOLD
summary_box.line.color.rgb = COLOR_GOLD
tf_sm = summary_box.text_frame
tf_sm.word_wrap = True

p_s1 = tf_sm.paragraphs[0]
p_s1.text = "Architectural Paradigm Manifesto"
p_s1.font.size = Pt(16)
p_s1.font.bold = True
p_s1.font.color.rgb = COLOR_NAVY
p_s1.space_after = Pt(6)

p_s2 = tf_sm.add_paragraph()
p_s2.text = "This Agentic AI DBA Platform transitions operational units from reactive disaster containment routines into a predictable data-driven environment. " \
           "By unifying continuous telemetry scans with localized advanced reasoning models, it delivers precision diagnostics without introducing infrastructure risk.\n\n" \
           "Design Philosophy:\n" \
           "AI-Assisted. Governance-Controlled. Human-Approved."
p_s2.font.size = Pt(13)
p_s2.font.color.rgb = COLOR_TEXT_DARK

# Bottom Technical Stack Badges Row
badges = ["Python", "FastAPI", "SQL Server", "OpenAI", "Azure-Ready", "MCP Architecture"]
for idx, b_text in enumerate(badges):
    bx = Inches(6.9) + idx * Inches(0.98)
    b_shape = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(5.7), Inches(0.9), Inches(0.4))
    b_shape.fill.solid()
    b_shape.fill.fore_color.rgb = COLOR_NAVY
    b_shape.line.fill.background()
    p_b = b_shape.text_frame.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    p_b.text = b_text
    p_b.font.size = Pt(9)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_WHITE

add_footer(slide12, 12)

# Save Document Assets Locally
prs.save("Autonomous_AI_DBA_Operations_Platform_Demo.pptx")
print("Presentation file 'Autonomous_AI_DBA_Operations_Platform_Demo.pptx' generated successfully.")